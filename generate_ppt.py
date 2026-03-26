from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Slide 1: Title
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "ClimaSense Intelligence Hub"
subtitle.text = "Automated Climate Stress & Anomaly Detection Protocol\nPresented by: Vansh Agrawal"

# Helper for bullet slides
def add_bullet_slide(title_text, bullets):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = title_text
    
    tf = body_shape.text_frame
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet

# Slide 2: The Problem
add_bullet_slide(
    "The Problem with Modern Climate Data",
    [
        "Existing weather applications are static, boring, and lack predictive depth.",
        "Climate anomalies (heatwaves, extreme pollution) are increasing.",
        "Tools to monitor anomalies visually remain unintuitive.",
        "Lack of deep integration between raw telemetry, geospatial visuals, and predictive Machine Learning."
    ]
)

# Slide 3: The Solution
add_bullet_slide(
    "A Paradigm Shift in Environmental Monitoring",
    [
        "ClimaSense: Transforms raw datasets (2024-2025) and live APIs into an immersive Situation Room dashboard.",
        "End-to-end data science pipeline visualized through a Pristine Light / Tech Dark UI.",
        "Actively detects extreme weather events using explainable integrated ML models.",
        "Provides predictive trajectories instead of just showing historical charts."
    ]
)

# Slide 4: System Architecture
add_bullet_slide(
    "System Architecture & Stack",
    [
        "Backend Services: FastAPI, Python, Scikit-Learn, XGBoost (aqi_classifier.pkl).",
        "Data Processing: Pandas-driven scripts (regenerate_forecast.py) for rolling predictions.",
        "Frontend UI: React 19, Vite, Tailwind CSS v3.",
        "Visualizations: Framer Motion (animations), Recharts (analytics), Leaflet + Esri World Imagery."
    ]
)

# Slide 5: Predictive Analytics
add_bullet_slide(
    "3-Year Forecasting & AI Inference",
    [
        "Trajectory Engine: ML-driven projections for 2026-2028 climate stress across 10 major Indian nodes.",
        "Quantum Inference Probe: On-demand XGBoost AI model evaluates real-time environment metrics (temp, humidity, wind).",
        "Backend resolves Live AQI severity classes dynamically.",
        "Interactive 5-year trajectory chart for high-level forecasting."
    ]
)

# Slide 6: Geospatial Tracking
add_bullet_slide(
    "Immersive Satellite HUD",
    [
        "High-resolution Esri Satellite maps targeting specific city coordinates seamlessly.",
        "Overlaid with Animated HUD 'pings' simulating real-time satellite node connections.",
        "Displays dynamic calibration alerts for localized thermal and pressure variances."
    ]
)

# Slide 7: Anomaly Detection
add_bullet_slide(
    "The Anomalies Dashboard",
    [
        "Dynamic Radar Spectrum Analysis evaluating Temperature, Pressure, Moisture, Radiation, and CO2.",
        "Live spatial anomaly logs decrypting raw atmospheric signals simulating deep diagnostics.",
        "Global sensor ticker providing real-time latency and status updates across the network."
    ]
)

# Slide 8: Mobile Integration
add_bullet_slide(
    "Mobile Integration & Execution",
    [
        "Fully responsive cross-device architecture.",
        "QR Code Connectivity: Hosted local-network bridging natively built into the app.",
        "Supervisors can scan and monitor the dashboard on mobile devices seamlessly.",
        "Report Generation: Extracts system state variables into downloadable text reports."
    ]
)

# Slide 9: Future Roadmap
add_bullet_slide(
    "Expanding the Array (Future Roadmap)",
    [
        "Integrating autonomous policy impact simulations.",
        "Expanding datasets to include more global nodes internationally.",
        "Advanced seasonal risk forecasting methodologies.",
        "Integration of deep learning for autonomous pattern recognition."
    ]
)

# Slide 10: Conclusion
add_bullet_slide(
    "Conclusion",
    [
        "The climate may be stressed, but the code isn't.",
        "Live Demo.",
        "Open for questions.",
        "GitHub: vansh070605"
    ]
)

prs.save('ClimaSense_Presentation.pptx')
